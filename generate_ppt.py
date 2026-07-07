import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide layout helper (6 is a blank layout)
    blank_layout = prs.slide_layouts[6]
    
    # Colors
    BG_COLOR = RGBColor(7, 8, 13)
    WHITE = RGBColor(255, 255, 255)
    MUTED = RGBColor(159, 164, 176)
    CYAN = RGBColor(0, 240, 255)
    GREEN = RGBColor(0, 255, 102)
    ORANGE = RGBColor(255, 136, 0)
    DARK_PANEL = RGBColor(24, 26, 37)

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_logo(slide):
        logo_path = os.path.join("images", "logo.jpg")
        if os.path.exists(logo_path):
            # Place the Sundar Innovations logo at the top-right corner
            slide.shapes.add_picture(logo_path, Inches(11.0), Inches(0.4), width=Inches(1.6), height=Inches(0.65))

    def add_title(slide, category, title_text, category_color):
        # Category tracking tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = 'Outfit'
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = category_color
        
        # Main slide title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Outfit'
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE
        
        # Add logo on standard slide
        add_logo(slide)

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    add_logo(slide1)
    
    # Title Text Frame (Left half)
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(7.5), Inches(4.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    # Main Product Name
    p1 = tf.paragraphs[0]
    p1.text = "SUNDAR CARBON"
    p1.font.name = 'Outfit'
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = GREEN
    p1.space_after = Pt(0)
    
    p2 = tf.add_paragraph()
    p2.text = "ANTIDOTE PLUS"
    p2.font.name = 'Outfit'
    p2.font.size = Pt(54)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(15)
    
    p3 = tf.add_paragraph()
    p3.text = "with Active Booster Dose (Ethanol Defense)"
    p3.font.name = 'Outfit'
    p3.font.size = Pt(22)
    p3.font.bold = True
    p3.font.color.rgb = CYAN
    p3.space_after = Pt(25)
    
    p4 = tf.add_paragraph()
    p4.text = "The ultimate engine rejuvenation, performance boost, and certified fuel system defense for BS4 & BS6 petrol vehicles."
    p4.font.name = 'Inter'
    p4.font.size = Pt(14)
    p4.font.color.rgb = MUTED
    
    # Add Product Image (Right half)
    img_path = os.path.join("images", "product_showcase.png")
    if os.path.exists(img_path):
        slide1.shapes.add_picture(img_path, Inches(8.5), Inches(1.2), width=Inches(4.0), height=Inches(5.0))

    # ==========================================
    # SLIDE 2: THE PROBLEM (ETHANOL THREAT)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_title(slide2, "The Threat", "The Silent Engine Killer", ORANGE)
    
    # Left Content Box (Bullets)
    left_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "Mandatory government ethanol-blending (up to 20%) in petrol creates high-risk scenarios for engines:"
    p.font.name = 'Inter'
    p.font.size = Pt(15)
    p.font.color.rgb = MUTED
    p.space_after = Pt(20)
    
    bullets = [
        ("Water Absorption (Hygroscopic):", " Ethanol actively absorbs moisture from air, creating a corrosive water-fuel mixture at the bottom of the fuel tank."),
        ("Corrosive Attack:", " Damages metal tanks, rusts internal fuel lines, and degrades rubber seals and hose couplings."),
        ("BS4 & BS6 Failure:", " High-pressure fuel injectors clog quickly due to separation deposits, leading to engine misfires, knocking, and power loss.")
    ]
    
    for title, text in bullets:
        p_b = tf_left.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = ORANGE
        
        # Add bold label
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        # Add normal text
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(14)
        
    # Right Side Graphic Panel (Visualizing damage)
    shape = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_PANEL
    shape.line.color.rgb = ORANGE
    shape.line.width = Pt(1.5)
    
    tf_panel = shape.text_frame
    tf_panel.word_wrap = True
    tf_panel.margin_left = Inches(0.4)
    tf_panel.margin_right = Inches(0.4)
    tf_panel.margin_top = Inches(0.4)
    
    p_p1 = tf_panel.paragraphs[0]
    p_p1.text = "UNPROTECTED FUEL SYSTEM"
    p_p1.font.name = 'Outfit'
    p_p1.font.bold = True
    p_p1.font.size = Pt(18)
    p_p1.font.color.rgb = ORANGE
    p_p1.space_after = Pt(15)
    
    panel_bullets = [
        "Phase separation isolates water layer at fuel tank base",
        "Corroding fuel pumps and rusting injectors",
        "Decreased thermal combustion efficiency",
        "Higher emissions & costly mechanical wear-and-tear"
    ]
    for b in panel_bullets:
        p_pb = tf_panel.add_paragraph()
        p_pb.text = "⚠️ " + b
        p_pb.font.name = 'Inter'
        p_pb.font.size = Pt(13)
        p_pb.font.color.rgb = WHITE
        p_pb.space_after = Pt(10)

    # ==========================================
    # SLIDE 3: THE SOLUTION (ACTIVE SHIELD)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_title(slide3, "The Shield", "Active Ethanol Defense", CYAN)
    
    # Left Content Box (Bullets)
    left_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left3 = left_box3.text_frame
    tf_left3.word_wrap = True
    
    p = tf_left3.paragraphs[0]
    p.text = "The Booster Dose is engineered to chemically safeguard your entire fuel line:"
    p.font.name = 'Inter'
    p.font.size = Pt(15)
    p.font.color.rgb = MUTED
    p.space_after = Pt(20)
    
    bullets3 = [
        ("Combats Phase Separation:", " Keeps ethanol and moisture fully dispersed and bonded within the fuel. Prevents water pooling at the tank base."),
        ("Micro-Coating Shield:", " Coats fuel pump parts, lines, and injectors with a molecular anti-corrosive layer that repels water molecules."),
        ("Nozzle De-Clogging:", " Keeps high-precision BS4 and BS6 fuel injector spray orifices completely clean and clear of deposits.")
    ]
    
    for title, text in bullets3:
        p_b = tf_left3.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = CYAN
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(14)
        
    # Right Side Graphic Panel (Visualizing protection)
    shape3 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.5))
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = DARK_PANEL
    shape3.line.color.rgb = CYAN
    shape3.line.width = Pt(1.5)
    
    tf_panel3 = shape3.text_frame
    tf_panel3.word_wrap = True
    tf_panel3.margin_left = Inches(0.4)
    tf_panel3.margin_right = Inches(0.4)
    tf_panel3.margin_top = Inches(0.4)
    
    p_p3 = tf_panel3.paragraphs[0]
    p_p3.text = "BOOSTER DOSE SHIELD ON"
    p_p3.font.name = 'Outfit'
    p_p3.font.bold = True
    p_p3.font.size = Pt(18)
    p_p3.font.color.rgb = CYAN
    p_p3.space_after = Pt(15)
    
    panel_bullets3 = [
        "Corrosion prevented on fuel lines, tanks, and valves",
        "Clean, stable fuel combustion under all blended states",
        "Maximum nozzle atomization (BS4 & BS6 optimized)",
        "Extended life for high-pressure fuel pumps"
    ]
    for b in panel_bullets3:
        p_pb = tf_panel3.add_paragraph()
        p_pb.text = "✔️ " + b
        p_pb.font.name = 'Inter'
        p_pb.font.size = Pt(13)
        p_pb.font.color.rgb = WHITE
        p_pb.space_after = Pt(10)

    # ==========================================
    # SLIDE 4: PERFORMANCE SURGE
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_title(slide4, "Performance Surge", "Certified Engine Output", GREEN)
    
    # Subtitle intro
    sub_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.0), Inches(0.5))
    sub_tf = sub_box.text_frame
    p_sub = sub_tf.paragraphs[0]
    p_sub.text = "Certified performance values representing significant operational efficiency improvements:"
    p_sub.font.name = 'Inter'
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = MUTED
    
    # 3 Stat Columns
    stats = [
        ("+30%", "FUEL EFFICIENCY", "Increases vehicle mileage and overall fuel economy by up to 30%.", GREEN),
        ("+80%", "TORQUE & POWER", "Delivers up to an 80% increase in engine pickup, power, and torque.", CYAN),
        ("-80%", "COOLER RUNNING", "Decreases external engine running temperatures by up to 80%.", ORANGE)
    ]
    
    for i, (val, title, desc, color) in enumerate(stats):
        # Stat Card Frame
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*4.0), Inches(2.3), Inches(3.7), Inches(4.0))
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_PANEL
        card.line.color.rgb = color
        card.line.width = Pt(1)
        
        card_tf = card.text_frame
        card_tf.word_wrap = True
        card_tf.margin_left = Inches(0.3)
        card_tf.margin_right = Inches(0.3)
        card_tf.margin_top = Inches(0.4)
        
        # Giant Number
        p_val = card_tf.paragraphs[0]
        p_val.text = val
        p_val.font.name = 'Outfit'
        p_val.font.size = Pt(54)
        p_val.font.bold = True
        p_val.font.color.rgb = color
        p_val.alignment = PP_ALIGN.CENTER
        p_val.space_after = Pt(5)
        
        # Label
        p_lbl = card_tf.add_paragraph()
        p_lbl.text = title
        p_lbl.font.name = 'Outfit'
        p_lbl.font.size = Pt(14)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = WHITE
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.space_after = Pt(15)
        
        # Description
        p_desc = card_tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Inter'
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = MUTED
        p_desc.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 5: MAINTENANCE & DURABILITY
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_title(slide5, "Maintenance & Durability", "Extreme Wear Reduction", CYAN)
    
    # Left Content
    left_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left5 = left_box5.text_frame
    tf_left5.word_wrap = True
    
    p_intro5 = tf_left5.paragraphs[0]
    p_intro5.text = "Minimizes engine downtime and cuts vehicle maintenance overhead:"
    p_intro5.font.name = 'Inter'
    p_intro5.font.size = Pt(15)
    p_intro5.font.color.rgb = MUTED
    p_intro5.space_after = Pt(20)
    
    bullets5 = [
        ("85% Less Wear & Tear:", " Decreases mechanical engine parts wear-and-tear by up to 85% by creating a resilient lubricant boundary layer."),
        ("5x Extended Oil Life:", " Reduces oil oxidation and soot contamination, extending oil longevity up to 5x."),
        ("5x Extended Engine Life:", " Protects internal components from micro-abrasions and combustion stress."),
        ("Refined Smoothness:", " Engine noise drops by 60% and vibrations cut by 40% for a smooth drive. Increases AC performance by up to 60% by reducing load.")
    ]
    
    for title, text in bullets5:
        p_b = tf_left5.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = CYAN
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(10)
        
    # Right Side Panel
    shape5 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.5))
    shape5.fill.solid()
    shape5.fill.fore_color.rgb = DARK_PANEL
    shape5.line.color.rgb = CYAN
    shape5.line.width = Pt(1)
    
    tf_panel5 = shape5.text_frame
    tf_panel5.word_wrap = True
    tf_panel5.margin_left = Inches(0.4)
    tf_panel5.margin_right = Inches(0.4)
    tf_panel5.margin_top = Inches(0.4)
    
    p_p5 = tf_panel5.paragraphs[0]
    p_p5.text = "FINANCIAL IMPACT"
    p_p5.font.name = 'Outfit'
    p_p5.font.bold = True
    p_p5.font.size = Pt(18)
    p_p5.font.color.rgb = CYAN
    p_p5.space_after = Pt(15)
    
    financial_points = [
        "Fewer engine oil changes (saves ₹3,000+ per change)",
        "Zero premature fuel system replacement costs",
        "Less downtime for fleet vehicles, maximizing revenue",
        "Up to ₹15,000+ in annual fuel & maintenance savings"
    ]
    for b in financial_points:
        p_pb = tf_panel5.add_paragraph()
        p_pb.text = "₹ " + b
        p_pb.font.name = 'Inter'
        p_pb.font.size = Pt(13)
        p_pb.font.color.rgb = WHITE
        p_pb.space_after = Pt(12)

    # ==========================================
    # SLIDE 6: ENVIRONMENTAL IMPACT
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_title(slide6, "Environmental Impact", "Clean Combustion", GREEN)
    
    # Left Content
    left_box6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left6 = left_box6.text_frame
    tf_left6.word_wrap = True
    
    p_intro6 = tf_left6.paragraphs[0]
    p_intro6.text = "A revolutionary leap in emission reductions and exhaust quality:"
    p_intro6.font.name = 'Inter'
    p_intro6.font.size = Pt(15)
    p_intro6.font.color.rgb = MUTED
    p_intro6.space_after = Pt(20)
    
    bullets6 = [
        ("99.9% NOx & SOx Elimination:", " Near-complete removal of highly toxic nitrogen oxides and sulfur oxides in the exhaust stream."),
        ("90% Arsenic Reduction:", " Drastically minimizes toxic arsenic particulate exhaust by up to 90%."),
        ("Oxygen Boost (+50% O₂):", " Dramatically increases clean, breathable oxygen levels in the tailpipe output."),
        ("Neutralizes Footprint:", " Optimizes engine combustion to neutralize standard carbon footprint outputs.")
    ]
    
    for title, text in bullets6:
        p_b = tf_left6.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = GREEN
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(10)
        
    # Right Side Graphic Panel (Visualizing clean tailpipe)
    shape6 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.5))
    shape6.fill.solid()
    shape6.fill.fore_color.rgb = DARK_PANEL
    shape6.line.color.rgb = GREEN
    shape6.line.width = Pt(1)
    
    tf_panel6 = shape6.text_frame
    tf_panel6.word_wrap = True
    tf_panel6.margin_left = Inches(0.4)
    tf_panel6.margin_right = Inches(0.4)
    tf_panel6.margin_top = Inches(0.4)
    
    p_p6 = tf_panel6.paragraphs[0]
    p_p6.text = "ECOLOGICAL COMPLIANCE"
    p_p6.font.name = 'Outfit'
    p_p6.font.bold = True
    p_p6.font.size = Pt(18)
    p_p6.font.color.rgb = GREEN
    p_p6.space_after = Pt(15)
    
    eco_points = [
        "Highly recommended for PUC emission compliance tests",
        "Converts exhaust elements into cleaner air compounds",
        "Neutralizes carbon output through extreme fuel efficiency",
        "Reduces toxic smog-forming compounds to near-zero levels"
    ]
    for b in eco_points:
        p_pb = tf_panel6.add_paragraph()
        p_pb.text = "🌱 " + b
        p_pb.font.name = 'Inter'
        p_pb.font.size = Pt(13)
        p_pb.font.color.rgb = WHITE
        p_pb.space_after = Pt(12)

    # ==========================================
    # SLIDE 7: PEACE OF MIND
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_title(slide7, "Peace of Mind", "Long-Term Reliability", GREEN)
    
    # Left Content
    left_box7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left7 = left_box7.text_frame
    tf_left7.word_wrap = True
    
    p_intro7 = tf_left7.paragraphs[0]
    p_intro7.text = "Engineered for maximum duration and passenger peace of mind:"
    p_intro7.font.name = 'Inter'
    p_intro7.font.size = Pt(15)
    p_intro7.font.color.rgb = MUTED
    p_intro7.space_after = Pt(25)
    
    p_guar = tf_left7.add_paragraph()
    p_guar.text = "1,00,000 KM ACTIVE DURATION"
    p_guar.font.name = 'Outfit'
    p_guar.font.size = Pt(20)
    p_guar.font.bold = True
    p_guar.font.color.rgb = GREEN
    p_guar.space_after = Pt(10)
    
    p_guar_desc = tf_left7.add_paragraph()
    p_guar_desc.text = "A single complete application remains fully active in the vehicle's engine and lubrication systems for 1,00,000 kilometers, requiring zero top-ups or frequent treatment renewals."
    p_guar_desc.font.name = 'Inter'
    p_guar_desc.font.size = Pt(14)
    p_guar_desc.font.color.rgb = MUTED
    p_guar_desc.space_after = Pt(25)
    
    p_note_header = tf_left7.add_paragraph()
    p_note_header.text = "PRESENTER NOTE:"
    p_note_header.font.name = 'Outfit'
    p_note_header.font.size = Pt(14)
    p_note_header.font.bold = True
    p_note_header.font.color.rgb = CYAN
    p_note_header.space_after = Pt(5)
    
    p_note_text = tf_left7.add_paragraph()
    p_note_text.text = "Keep presentations focused on the long-term cost benefits, mechanical reliability, and peace of mind this 1,00,000 km timeline offers to individual car owners and corporate fleet managers."
    p_note_text.font.name = 'Inter'
    p_note_text.font.size = Pt(13)
    p_note_text.font.color.rgb = MUTED

    # Right Side Graphic Panel (Seal representation)
    shape7 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.5))
    shape7.fill.solid()
    shape7.fill.fore_color.rgb = DARK_PANEL
    shape7.line.color.rgb = GREEN
    shape7.line.width = Pt(1.5)
    
    tf_panel7 = shape7.text_frame
    tf_panel7.word_wrap = True
    tf_panel7.margin_left = Inches(0.4)
    tf_panel7.margin_right = Inches(0.4)
    tf_panel7.margin_top = Inches(0.8)
    
    p_p7 = tf_panel7.paragraphs[0]
    p_p7.text = "1,00,000 KM"
    p_p7.font.name = 'Outfit'
    p_p7.font.bold = True
    p_p7.font.size = Pt(40)
    p_p7.font.color.rgb = GREEN
    p_p7.alignment = PP_ALIGN.CENTER
    p_p7.space_after = Pt(10)
    
    p_p7_2 = tf_panel7.add_paragraph()
    p_p7_2.text = "ACTIVE EFFICACY GUARANTEE"
    p_p7_2.font.name = 'Outfit'
    p_p7_2.font.bold = True
    p_p7_2.font.size = Pt(18)
    p_p7_2.font.color.rgb = WHITE
    p_p7_2.alignment = PP_ALIGN.CENTER
    p_p7_2.space_after = Pt(20)
    
    p_p7_3 = tf_panel7.add_paragraph()
    p_p7_3.text = "Certified to protect high-pressure fuel injectors and mechanical cylinders across a massive driving range."
    p_p7_3.font.name = 'Inter'
    p_p7_3.font.size = Pt(13)
    p_p7_3.font.color.rgb = MUTED
    p_p7_3.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_filename = "sundar_carbon_antidote_plus_with_logo.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as {output_filename}")

if __name__ == "__main__":
    create_deck()
