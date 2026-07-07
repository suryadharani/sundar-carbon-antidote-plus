import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide layout helper (6 is a blank layout)
    blank_layout = prs.slide_layouts[6]
    
    # Colors
    BG_COLOR = RGBColor(7, 8, 13)
    WHITE = RGBColor(255, 255, 255)
    MUTED = RGBColor(159, 164, 176)
    CYAN = RGBColor(0, 240, 255)
    GREEN = RGBColor(0, 255, 102)
    ORANGE = RGBColor(255, 136, 0)
    DARK_PANEL = RGBColor(24, 26, 37)

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_logo(slide):
        logo_path = os.path.join("images", "logo.jpg")
        if os.path.exists(logo_path):
            # Place the Sundar Innovations logo at the top-right corner
            slide.shapes.add_picture(logo_path, Inches(11.0), Inches(0.4), width=Inches(1.6), height=Inches(0.65))

    def add_title(slide, category, title_text, category_color):
        # Category tracking tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = 'Outfit'
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = category_color
        
        # Main slide title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Outfit'
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE
        
        # Add logo on standard slide
        add_logo(slide)

    def add_slide_picture(slide, img_name, border_color):
        # Draw the frame shape
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_PANEL
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        
        # Overlay the picture inside the frame with padding
        img_path = os.path.join("images", img_name)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(7.6), Inches(1.9), width=Inches(4.8), height=Inches(4.3))

    # ==========================================
    # SLIDE 1: TITLE SLIDE (India Edition)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    add_logo(slide1)
    
    # Title Text Frame (Left half)
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.5), Inches(4.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    # Hindi Local Banner Subtitle
    p0 = tf.paragraphs[0]
    p0.text = "भारत के वाहनों के लिए रक्षा कवच"
    p0.font.name = 'Mangal' # Standard Windows Devanagari/Hindi font
    p0.font.size = Pt(22)
    p0.font.bold = True
    p0.font.color.rgb = GREEN
    p0.space_after = Pt(5)
    
    # Main Product Name
    p1 = tf.add_paragraph()
    p1.text = "SUNDAR CARBON"
    p1.font.name = 'Outfit'
    p1.font.size = Pt(52)
    p1.font.bold = True
    p1.font.color.rgb = GREEN
    p1.space_after = Pt(0)
    
    p2 = tf.add_paragraph()
    p2.text = "ANTIDOTE PLUS"
    p2.font.name = 'Outfit'
    p2.font.size = Pt(52)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(12)
    
    p3 = tf.add_paragraph()
    p3.text = "with Active Booster Dose (Ethanol Defense)"
    p3.font.name = 'Outfit'
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = CYAN
    p3.space_after = Pt(20)
    
    p4 = tf.add_paragraph()
    p4.text = "The ultimate engine rejuvenation, performance boost, and fuel system defense calibrated for all Petrol & Diesel vehicles on Indian roads."
    p4.font.name = 'Inter'
    p4.font.size = Pt(13)
    p4.font.color.rgb = MUTED
    
    # Add Product Image (Right half)
    img_path = os.path.join("images", "product_showcase.png")
    if os.path.exists(img_path):
        slide1.shapes.add_picture(img_path, Inches(8.6), Inches(1.5), width=Inches(3.8), height=Inches(4.8))

    # ==========================================
    # SLIDE 2: THE PROBLEM (NATIONAL E20 MANDATE)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_title(slide2, "The Fuel Threat", "The National E20 Challenge", ORANGE)
    
    # Left Content Box (Bullets)
    left_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "India's mandatory transition to E20 (20% Ethanol) blended petrol poses severe nationwide fuel system threats:"
    p.font.name = 'Inter'
    p.font.size = Pt(15)
    p.font.color.rgb = MUTED
    p.space_after = Pt(20)
    
    bullets = [
        ("National Biofuel Policy Mandate:", " India has accelerated E20 fuel availability nationwide to reduce oil import bills, making ethanol-blended petrol standard at all pumps."),
        ("Humidity & Water Phase Separation:", " Monsoon seasons and high humidity across India cause ethanol to rapidly absorb moisture. This triggers phase separation, pooling corrosive water inside the vehicle fuel tank."),
        ("BS4 & BS6 Injector Risk:", " Modern high-pressure fuel injectors clog quickly with gummy deposits, leading to engine knocking, lost mileage, and misfires.")
    ]
    
    for title, text in bullets:
        p_b = tf_left.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = ORANGE
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(12)
        
    add_slide_picture(slide2, "ethanol_damage.png", ORANGE)

    # ==========================================
    # SLIDE 3: THE SHIELD (BOOSTER DOSE)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_title(slide3, "The Shield", "Booster Dose (Ethanol Defense)", CYAN)
    
    # Left Content Box (Bullets)
    left_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left3 = left_box3.text_frame
    tf_left3.word_wrap = True
    
    p = tf_left3.paragraphs[0]
    p.text = "The active Booster Dose provides an immediate chemical defense shield for petrol engines:"
    p.font.name = 'Inter'
    p.font.size = Pt(15)
    p.font.color.rgb = MUTED
    p.space_after = Pt(20)
    
    bullets3 = [
        ("Binds Water & Prevents Pools:", " Keeps ethanol and water fully bonded in the fuel, ensuring it burns off smoothly without pooling in the tank."),
        ("Molecular Anti-Corrosive Layer:", " Coats fuel pump parts, lines, and injectors with a molecular anti-corrosive layer that repels acidic moisture."),
        ("Restores Injector Spray:", " Disperses blockages and keeps high-precision BS4 & BS6 injector orifices completely clean for fuel compliance.")
    ]
    
    for title, text in bullets3:
        p_b = tf_left3.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = CYAN
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(12)
        
    add_slide_picture(slide3, "molecular_shield.png", CYAN)

    # ==========================================
    # SLIDE 4: PERFORMANCE SURGE (INDIA'S LANDSCAPES)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_title(slide4, "Performance Surge", "Conquering India's Diverse Landscapes", GREEN)
    
    # Left Content Box (Bullets)
    left_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left4 = left_box4.text_frame
    tf_left4.word_wrap = True
    
    p = tf_left4.paragraphs[0]
    p.text = "Certified performance values representing major operational improvements on local routes:"
    p.font.name = 'Inter'
    p.font.size = Pt(15)
    p.font.color.rgb = MUTED
    p.space_after = Pt(20)
    
    bullets4 = [
        ("+30% Mileage Boost:", " Optimizes fuel combustion to increase mileage and fuel economy by up to 30%. Saves running costs on long highway runs."),
        ("+80% Torque & Power:", " Delivers an 80% increase in power and pickup. Critical for steep Ghat routes, high-altitude climbs, and overtaking maneuvers."),
        ("-80% Cooler Running:", " Decreases external engine running temperatures by up to 80%. Prevents overheating in dense metropolitan bumper-to-bumper traffic gridlocks (Delhi NCR, Mumbai, Bengaluru).")
    ]
    
    for title, text in bullets4:
        p_b = tf_left4.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = GREEN
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(14)
        
    add_slide_picture(slide4, "performance_surge.png", GREEN)

    # ==========================================
    # SLIDE 5: MAINTENANCE & DURABILITY
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_title(slide5, "Maintenance & Durability", "Extreme Wear & Cost Reduction", CYAN)
    
    # Left Content
    left_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left5 = left_box5.text_frame
    tf_left5.word_wrap = True
    
    p_intro5 = tf_left5.paragraphs[0]
    p_intro5.text = "Minimizes engine downtime and cuts annual maintenance budgets:"
    p_intro5.font.name = 'Inter'
    p_intro5.font.size = Pt(15)
    p_intro5.font.color.rgb = MUTED
    p_intro5.space_after = Pt(20)
    
    bullets5 = [
        ("85% Less Wear & Tear:", " Decreases mechanical engine parts wear by up to 85% by creating a resilient lubricant boundary layer."),
        ("5x Extended Oil Life:", " Reduces oil degradation and contamination, extending oil changes by 5 times. Saves ₹15,000+ annually for typical drivers."),
        ("5x Extended Engine Life:", " Protects internal pistons and valves from micro-abrasions and heat stress."),
        ("Refined Smoothness:", " Drops engine noise by 60% and vibrations by 40% on rough roads. Enhances vehicle AC performance by up to 60% by easing engine strain.")
    ]
    
    for title, text in bullets5:
        p_b = tf_left5.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = CYAN
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(10)
        
    add_slide_picture(slide5, "piston_wear.png", CYAN)

    # ==========================================
    # SLIDE 6: MISSION LIFE & CLEAN AIR (NCAP)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_title(slide6, "Eco Impact", "Mission LiFE & Clean Air Compliance", GREEN)
    
    # Left Content
    left_box6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left6 = left_box6.text_frame
    tf_left6.word_wrap = True
    
    p_intro6 = tf_left6.paragraphs[0]
    p_intro6.text = "Directly aligns with India's national eco initiatives and clean air targets:"
    p_intro6.font.name = 'Inter'
    p_intro6.font.size = Pt(15)
    p_intro6.font.color.rgb = MUTED
    p_intro6.space_after = Pt(20)
    
    bullets6 = [
        ("Mission LiFE (Lifestyle for Environment):", " PM Modi's national movement encouraging individual action to conserve fuel, optimize energy, and adopt green driving habits."),
        ("National Clean Air Programme (NCAP):", " Supports India's NCAP targets to reduce urban smog and particulate matter in major cities."),
        ("99.9% NOx & SOx Elimination:", " Near-complete removal of toxic nitrogen oxides and sulfur oxides in the exhaust stream."),
        ("90% Arsenic & +50% O₂ Boost:", " Drastically minimizes heavy metal arsenic emissions by 90% while adding 50% clean oxygen to tailpipe output.")
    ]
    
    for title, text in bullets6:
        p_b = tf_left6.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = GREEN
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(10)
        
    add_slide_picture(slide6, "clean_emissions.png", GREEN)

    # ==========================================
    # SLIDE 7: PEACE OF MIND
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_title(slide7, "Peace of Mind", "1,00,000 KM Efficacy Guarantee", GREEN)
    
    # Left Content
    left_box7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left7 = left_box7.text_frame
    tf_left7.word_wrap = True
    
    p_intro7 = tf_left7.paragraphs[0]
    p_intro7.text = "Engineered for maximum duration and passenger peace of mind:"
    p_intro7.font.name = 'Inter'
    p_intro7.font.size = Pt(15)
    p_intro7.font.color.rgb = MUTED
    p_intro7.space_after = Pt(25)
    
    p_guar = tf_left7.add_paragraph()
    p_guar.text = "भारत के गौरवशाली पथ के लिए दीर्घकालिक सुरक्षा"
    p_guar.font.name = 'Mangal'
    p_guar.font.size = Pt(18)
    p_guar.font.bold = True
    p_guar.font.color.rgb = GREEN
    p_guar.space_after = Pt(10)
    
    p_guar_desc = tf_left7.add_paragraph()
    p_guar_desc.text = "A single complete application remains fully active in the vehicle's engine and lubrication systems for 1,00,000 kilometers, requiring zero top-ups or frequent treatment renewals."
    p_guar_desc.font.name = 'Inter'
    p_guar_desc.font.size = Pt(14)
    p_guar_desc.font.color.rgb = MUTED
    p_guar_desc.space_after = Pt(20)
    
    p_note_header = tf_left7.add_paragraph()
    p_note_header.text = "PRESENTER NOTE:"
    p_note_header.font.name = 'Outfit'
    p_note_header.font.size = Pt(14)
    p_note_header.font.bold = True
    p_note_header.font.color.rgb = CYAN
    p_note_header.space_after = Pt(5)
    
    p_note_text = tf_left7.add_paragraph()
    p_note_text.text = "Keep presentations focused on the long-term cost benefits, mechanical reliability, and peace of mind this 1,00,000 km timeline offers to individual car owners and corporate fleet managers across India."
    p_note_text.font.name = 'Inter'
    p_note_text.font.size = Pt(13)
    p_note_text.font.color.rgb = MUTED
        
    add_slide_picture(slide7, "longevity_road.png", GREEN)

    # ==========================================
    # SLIDE 8: UNIVERSAL COMPATIBILITY
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_title(slide8, "Compatibility", "Universal Fuel Adaptability", CYAN)
    
    # Left Content
    left_box8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left8 = left_box8.text_frame
    tf_left8.word_wrap = True
    
    p_intro8 = tf_left8.paragraphs[0]
    p_intro8.text = "Formulated to optimize all standard combustible fuel engine configurations:"
    p_intro8.font.name = 'Inter'
    p_intro8.font.size = Pt(15)
    p_intro8.font.color.rgb = MUTED
    p_intro8.space_after = Pt(25)
    
    bullets8 = [
        ("Petrol Engines (Ethanol Defense):", " Actively guards delicate BS4 & BS6 injectors from corrosive sugarcane-blended E20 fuel deposits."),
        ("Diesel Engines (Heavy Duty):", " Prevents cylinder scaling, exhaust soot buildup, and injectors carbonization under heavy-duty cargo demands."),
        ("Multi-Fuel Performance:", " Works identically on petrol-powered commuter vehicles and diesel-powered cargo fleets, requiring no formula alterations.")
    ]
    
    for title, text in bullets8:
        p_b = tf_left8.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = CYAN
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(14)
        
    add_slide_picture(slide8, "universal_vehicles.png", CYAN)

    # Save presentation
    output_filename = "sundar_carbon_antidote_plus_india.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as {output_filename}")

if __name__ == "__main__":
    create_deck()
