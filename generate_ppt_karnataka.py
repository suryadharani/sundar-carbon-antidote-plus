import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide layout helper (6 is a blank layout)
    blank_layout = prs.slide_layouts[6]
    
    # Colors
    BG_COLOR = RGBColor(7, 8, 13)
    WHITE = RGBColor(255, 255, 255)
    MUTED = RGBColor(159, 164, 176)
    CYAN = RGBColor(0, 240, 255)
    GREEN = RGBColor(0, 255, 102)
    ORANGE = RGBColor(255, 136, 0)
    DARK_PANEL = RGBColor(24, 26, 37)

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_logo(slide):
        logo_path = os.path.join("images", "logo_transparent.png")
        if os.path.exists(logo_path):
            # Place the Sundar Innovations logo at the top-right corner
            slide.shapes.add_picture(logo_path, Inches(11.0), Inches(0.4), width=Inches(1.6), height=Inches(0.65))

    def add_title(slide, category, title_text, category_color):
        # Category tracking tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = 'Outfit'
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = category_color
        
        # Main slide title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Outfit'
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE
        
        # Add logo on standard slide
        add_logo(slide)

    def add_slide_picture(slide, img_name, border_color):
        # Insert the pre-bordered, center-cropped, and rounded image directly
        base_name = os.path.splitext(img_name)[0]
        bordered_img_name = f"{base_name}_bordered.png"
        img_path = os.path.join("images", bordered_img_name)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.8), width=Inches(5.0), height=Inches(4.5))

    # ==========================================
    # SLIDE 1: TITLE SLIDE (Karnataka Edition)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    add_logo(slide1)
    
    # Title Text Frame (Left half)
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.5), Inches(4.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    # Kannada Local Banner Subtitle
    p0 = tf.paragraphs[0]
    p0.text = "ಕನ್ನಡಿಗರ ವಾಹನಗಳಿಗೆ ಸುರಕ್ಷಾ ಕವಚ"
    p0.font.name = 'Tunga'
    p0.font.size = Pt(24)
    p0.font.bold = True
    p0.font.color.rgb = GREEN
    p0.space_after = Pt(5)
    
    # Main Product Name
    p1 = tf.add_paragraph()
    p1.text = "SUNDAR CARBON"
    p1.font.name = 'Outfit'
    p1.font.size = Pt(52)
    p1.font.bold = True
    p1.font.color.rgb = GREEN
    p1.space_after = Pt(0)
    
    p2 = tf.add_paragraph()
    p2.text = "ANTIDOTE PLUS®"
    p2.font.name = 'Outfit'
    p2.font.size = Pt(52)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(12)
    
    p3 = tf.add_paragraph()
    p3.text = "with Active Booster Dose (Ethanol Defense)"
    p3.font.name = 'Outfit'
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = CYAN
    p3.space_after = Pt(20)
    
    p4 = tf.add_paragraph()
    p4.text = "The ultimate engine rejuvenation, performance boost, and fuel system defense calibrated for all Petrol & Diesel vehicles on Karnataka's roads."
    p4.font.name = 'Inter'
    p4.font.size = Pt(13)
    p4.font.color.rgb = MUTED
    
    # Add Product Image (Right half)
    img_path = os.path.join("images", "product_showcase_glow.png")
    if os.path.exists(img_path):
        slide1.shapes.add_picture(img_path, Inches(8.6), Inches(1.5), width=Inches(3.8), height=Inches(4.8))

    # ==========================================
    # SLIDE 2: AWARENESS - THE POLICY CHALLENGE (Karnataka Edition)
    # ==========================================
    slide_awareness = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_awareness)
    add_title(slide_awareness, "Policy & Infrastructure", "ಇಂಧನ ಪರಿವರ್ತನೆಯ ಸವಾಲು (The Energy Transition Challenge)", ORANGE)
    
    left_box_aw = slide_awareness.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left_aw = left_box_aw.text_frame
    tf_left_aw.word_wrap = True
    
    p = tf_left_aw.paragraphs[0]
    p.text = "Karnataka's aggressive push for E20 biofuel (derived from regional sugarcane molasses) supports local agriculture but raises significant transit risks:"
    p.font.name = 'Inter'
    p.font.size = Pt(14.5)
    p.font.color.rgb = MUTED
    p.space_after = Pt(15)
    
    bullets_aw = [
        ("Sugarcane Molasses Scaling:", " Karnataka's high sugarcane output accelerates E20 blending, but ethanol solvent degradation impacts vehicle fuel systems directly."),
        ("Regional Depots & Humidity:", " Depots in coastal Karnataka (Mangaluru, Karwar) experience high relative humidity, causing phase separation and water accumulation in depot storage tanks."),
        ("Transit Fleet Downtime:", " Public and state bus fleets suffer from injector corrosion, leading to unscheduled repairs and route downtime.")
    ]
    
    for title, text in bullets_aw:
        p_b = tf_left_aw.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = ORANGE
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(6)
        
    add_slide_picture(slide_awareness, "energy_transition_challenge.png", ORANGE)

    # ==========================================
    # SLIDE 3: CONSIDERATION - INFRASTRUCTURE DEFENSE (Karnataka Edition)
    # ==========================================
    slide_consideration = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_consideration)
    add_title(slide_consideration, "Operational Economics", "ಸಾರ್ವಜನಿಕ ಆಸ್ತಿ ಮತ್ತು ಬಜೆಟ್ ಸಂರಕ್ಷಣೆ (Protecting Public Assets & Budgets)", CYAN)
    
    left_box_con = slide_consideration.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left_con = left_box_con.text_frame
    tf_left_con.word_wrap = True
    
    p = tf_left_con.paragraphs[0]
    p.text = "Managing regional public transportation and municipal fleets under the E20 mandate requires comparing key financial strategies:"
    p.font.name = 'Inter'
    p.font.size = Pt(14.5)
    p.font.color.rgb = MUTED
    p.space_after = Pt(15)
    
    bullets_con = [
        ("Premium Fuel Scaling Cost:", " Upgrading KSRTC and municipal fleets to commercial premium fuels is budgetarily impossible, compounding state transit deficits."),
        ("Component Failure Risks:", " Leaving fleets untreated results in premature injector failure, tank rust remediation, and lost operational hours."),
        ("Molecular Treatment Savings:", " Implementing a regular, low-cost molecular fuel treatment shields public assets without budget inflation.")
    ]
    
    for title, text in bullets_con:
        p_b = tf_left_con.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = CYAN
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(6)
        
    add_slide_picture(slide_consideration, "infrastructure_protection.png", CYAN)

    # ==========================================
    # SLIDE 4: THE THREAT (E20 FUEL CHALLENGE)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_title(slide2, "The Threat", "The Sugar State Fuel Challenge", ORANGE)
    
    # Left Content Box (Bullets)
    left_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "Mandatory E20 (20% Ethanol) blended petrol creates severe operational risks for vehicles across Karnataka:"
    p.font.name = 'Inter'
    p.font.size = Pt(15)
    p.font.color.rgb = MUTED
    p.space_after = Pt(20)
    
    bullets = [
        ("Sugarcane Mandate:", " As a major crop producer, Karnataka actively pushes ethanol fuels derived from sugarcane molasses."),
        ("Water Phase Separation:", " In humid climates like coastal Mangaluru, Udupi, and Karwar, ethanol absorbs moisture, causing water pooling at the tank bottom."),
        ("BS4 & BS6 Clogging:", " High-pressure fuel injectors clog rapidly with deposits, leading to misfires, engine knocking, and loss of pickup on local roads."),
        ("Isobutanol-Diesel Transition:", " Next-generation isobutanol-diesel blends degrade standard elastomers (O-rings) and reduce fuel lubricity, accelerating wear in high-pressure diesel injector systems.")
    ]
    
    for title, text in bullets:
        p_b = tf_left.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(11.5)
        p_b.font.color.rgb = ORANGE
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(6)
        
    add_slide_picture(slide2, "ethanol_damage.png", ORANGE)

    # ==========================================
    # SLIDE 3: THE SHIELD (BOOSTER DOSE)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_title(slide3, "The Shield", "Booster Dose (Ethanol & Engine Protection)", CYAN)
    
    # Left Content Box (Bullets)
    left_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left3 = left_box3.text_frame
    tf_left3.word_wrap = True
    
    p = tf_left3.paragraphs[0]
    p.text = "The active Booster Dose provides an immediate chemical defense shield for petrol engines:"
    p.font.name = 'Inter'
    p.font.size = Pt(15)
    p.font.color.rgb = MUTED
    p.space_after = Pt(15)
    
    bullets3 = [
        ("Combats Phase Separation:", " Helps maintain a stable ethanol–petrol blend by reducing moisture-related separation."),
        ("Active Corrosion Protection:", " Helps reduce moisture-driven corrosive activity associated with ethanol-blended petrol."),
        ("Supports Precision Fuel Injection:", " Helps maintain consistent fuel delivery and injector performance in modern BS4 & BS6 petrol engines."),
        ("Reduces Carbon Deposits:", " Helps remove carbon deposits from the EGR system and other critical engine components, supporting reliable sensor operation and reducing maintenance requirements.")
    ]
    
    for title, text in bullets3:
        p_b = tf_left3.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = CYAN
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(8)
        
    add_slide_picture(slide3, "molecular_shield.png", CYAN)

    # ==========================================
    # SLIDE 4: THE MISSING LINK (THERMODYNAMIC BREAKTHROUGH)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_title(slide4, "Technology", "Cleaner Combustion – The Missing Link in ICE", GREEN)
    
    # Left Content Box (Paragraphs)
    left_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left4 = left_box4.text_frame
    tf_left4.word_wrap = True
    
    p_ice1 = tf_left4.paragraphs[0]
    p_ice1.text = "The Certification Gap in Current Frameworks"
    p_ice1.font.name = 'Outfit'
    p_ice1.font.size = Pt(18)
    p_ice1.font.bold = True
    p_ice1.font.color.rgb = GREEN
    p_ice1.space_after = Pt(10)
    
    p_ice2 = tf_left4.add_paragraph()
    p_ice2.text = "Current global vehicle homologation frameworks—including CMVR (India), UNECE Regulations (Europe), U.S. EPA/NHTSA, Japan's Top Runner Program, and China's CAFC regulations—certify vehicles based on emissions, fuel economy, CO₂, safety, and durability. They do not prescribe minimum combustion efficiency or brake thermal efficiency as independent certification criteria. Sundar Carbon directly addresses this underlying energy conversion process, unlocking performance improvements beyond conventional outcome-based optimization."
    p_ice2.font.name = 'Inter'
    p_ice2.font.size = Pt(12)
    p_ice2.font.color.rgb = MUTED
    p_ice2.space_after = Pt(20)
    
    p_ice3 = tf_left4.add_paragraph()
    p_ice3.text = "Improving Energy Conversion at Molecular Level"
    p_ice3.font.name = 'Outfit'
    p_ice3.font.size = Pt(18)
    p_ice3.font.bold = True
    p_ice3.font.color.rgb = GREEN
    p_ice3.space_after = Pt(10)
    
    p_ice4 = tf_left4.add_paragraph()
    p_ice4.text = "Unlike conventional approaches that primarily optimize vehicles to meet regulatory outcomes, our technology improves the combustion process itself at the molecular level, enabling substantially greater conversion of fuel energy into useful work. This represents a breakthrough in fuel energy utilization by improving Energy Conversion Efficiency from ~25% to ~65%, converting significantly more fuel into useful engine power while reducing unburnt hydrocarbons, emissions, fuel consumption, and fuel wastage."
    p_ice4.font.name = 'Inter'
    p_ice4.font.size = Pt(12)
    p_ice4.font.color.rgb = MUTED

    # Right Side Graphic Panel (Visualizing the conversion jump)
    shape4 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.5))
    shape4.fill.solid()
    shape4.fill.fore_color.rgb = DARK_PANEL
    shape4.line.color.rgb = GREEN
    shape4.line.width = Pt(1.5)
    
    tf_panel4 = shape4.text_frame
    tf_panel4.word_wrap = True
    tf_panel4.margin_left = Inches(0.4)
    tf_panel4.margin_right = Inches(0.4)
    tf_panel4.margin_top = Inches(0.6)
    
    p_p4 = tf_panel4.paragraphs[0]
    p_p4.text = "ENERGY CONVERSION IMPACT"
    p_p4.font.name = 'Outfit'
    p_p4.font.bold = True
    p_p4.font.size = Pt(18)
    p_p4.font.color.rgb = GREEN
    p_p4.space_after = Pt(25)
    
    p_std = tf_panel4.add_paragraph()
    p_std.text = "Standard Engine:"
    p_std.font.name = 'Inter'
    p_std.font.size = Pt(14)
    p_std.font.color.rgb = WHITE
    p_std.space_after = Pt(2)
    
    p_std_val = tf_panel4.add_paragraph()
    p_std_val.text = "~25% Efficiency"
    p_std_val.font.name = 'Outfit'
    p_std_val.font.size = Pt(28)
    p_std_val.font.bold = True
    p_std_val.font.color.rgb = ORANGE
    p_std_val.space_after = Pt(20)
    
    p_sun = tf_panel4.add_paragraph()
    p_sun.text = "With Sundar Carbon Antidotplus® + Booster Dose :"
    p_sun.font.name = 'Inter'
    p_sun.font.size = Pt(14)
    p_sun.font.color.rgb = WHITE
    p_sun.space_after = Pt(2)
    
    p_sun_val = tf_panel4.add_paragraph()
    p_sun_val.text = "~65% Efficiency"
    p_sun_val.font.name = 'Outfit'
    p_sun_val.font.size = Pt(36)
    p_sun_val.font.bold = True
    p_sun_val.font.color.rgb = GREEN
    p_sun_val.space_after = Pt(20)
    
    p_footer = tf_panel4.add_paragraph()
    p_footer.text = "Direct molecular combustion optimization reduces fuel wastage at source."
    p_footer.font.name = 'Inter'
    p_footer.font.size = Pt(11)
    p_footer.font.color.rgb = MUTED

    # ==========================================
    # SLIDE 5: PERFORMANCE SURGE
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_title(slide5, "Performance Surge", "Conquering Karnataka's Terrains", GREEN)
    
    # Left Content Box (Bullets)
    left_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left5 = left_box5.text_frame
    tf_left5.word_wrap = True
    
    p = tf_left5.paragraphs[0]
    p.text = "Certified performance values representing major operational improvements on local routes:"
    p.font.name = 'Inter'
    p.font.size = Pt(15)
    p.font.color.rgb = MUTED
    p.space_after = Pt(20)
    
    bullets5 = [
        ("+30% Mileage Boost:", " Optimizes fuel combustion to increase mileage and fuel economy by up to 30%. Saves running costs on highway trips."),
        ("+80% Torque & Power:", " Delivers an 80% increase in power and pickup. Critical for steep Western Ghats climbs (Charmadi, Shiradi) and Nandi Hills."),
        ("-80% Cooler Running:", " Decreases external engine running temperatures by up to 80%. Prevents overheating in Bengaluru's bumper-to-bumper gridlock.")
    ]
    
    for title, text in bullets5:
        p_b = tf_left5.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = GREEN
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(14)
        
    add_slide_picture(slide5, "performance_surge.png", GREEN)

    # ==========================================
    # SLIDE 6: MAINTENANCE & DURABILITY
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_title(slide6, "Maintenance & Durability", "Extreme Wear & Cost Reduction", CYAN)
    
    # Left Content
    left_box6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left6 = left_box6.text_frame
    tf_left6.word_wrap = True
    
    p_intro6 = tf_left6.paragraphs[0]
    p_intro6.text = "Minimizes engine downtime and cuts annual maintenance budgets:"
    p_intro6.font.name = 'Inter'
    p_intro6.font.size = Pt(15)
    p_intro6.font.color.rgb = MUTED
    p_intro6.space_after = Pt(20)
    
    bullets6 = [
        ("85% Less Wear & Tear:", " Decreases mechanical engine parts wear by up to 85% by creating a resilient lubricant boundary layer."),
        ("5x Extended Oil Life:", " Reduces oil degradation and contamination, extending oil changes by 5 times. Saves ₹15,000+ annually for typical drivers."),
        ("5x Extended Engine Life:", " Protects internal pistons and valves from micro-abrasions and heat stress."),
        ("50% Improved Gearbox Life:", " Reduces transmission wear and friction, enhancing gear operation and extending gearbox service life by up to 50%."),
        ("20% Longer Tyre Life:", " Promotes smoother power delivery and reduced drivetrain stress, contributing to up to 20% longer tyre life."),
        ("Refined Smoothness:", " Drops engine noise by 60% and vibrations by 40% on rough roads. Enhances vehicle AC performance by up to 60% by easing engine strain.")
    ]
    
    for title, text in bullets6:
        p_b = tf_left6.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = CYAN
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(6)
        
    add_slide_picture(slide6, "drivetrain_wear.png", CYAN)

    # ==========================================
    # SLIDE 7: ENVIRONMENTAL IMPACT
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_title(slide7, "Environmental Impact", "Green Karnataka Initiative", GREEN)
    
    # Left Content
    left_box7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left7 = left_box7.text_frame
    tf_left7.word_wrap = True
    
    p_intro7 = tf_left7.paragraphs[0]
    p_intro7.text = "Supports clean air initiatives in cities like Bengaluru, Mysuru, and Hubballi-Dharwad:"
    p_intro7.font.name = 'Inter'
    p_intro7.font.size = Pt(15)
    p_intro7.font.color.rgb = MUTED
    p_intro7.space_after = Pt(20)
    
    bullets7 = [
        ("99.9% NOx & SOx Elimination:", " Near-complete removal of toxic nitrogen oxides and sulfur oxides in the exhaust stream."),
        ("90% Arsenic Reduction:", " Drastically minimizes toxic heavy-metal arsenic particulate exhaust by up to 90%."),
        ("Oxygen Boost (+50% O₂):", " Dramatically increases clean, breathable oxygen levels in the tailpipe output."),
        ("Combustion Optimization:", " Optimizes engine combustion to neutralize standard carbon footprint outputs and ensure easy PUC compliance.")
    ]
    
    for title, text in bullets7:
        p_b = tf_left7.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = GREEN
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(10)
        
    add_slide_picture(slide7, "clean_emissions.png", GREEN)

    # ==========================================
    # SLIDE 8: PEACE OF MIND
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_title(slide8, "Peace of Mind", "1,00,000 KM Efficacy Guarantee", GREEN)
    
    # Left Content
    left_box8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left8 = left_box8.text_frame
    tf_left8.word_wrap = True
    
    p_intro8 = tf_left8.paragraphs[0]
    p_intro8.text = "Engineered for maximum duration and passenger peace of mind:"
    p_intro8.font.name = 'Inter'
    p_intro8.font.size = Pt(15)
    p_intro8.font.color.rgb = MUTED
    p_intro8.space_after = Pt(25)
    
    p_guar = tf_left8.add_paragraph()
    p_guar.text = "ಒಮ್ಮೆ ಬಳಸಿ, ದೀರ್ಘಕಾಲದವರೆಗೆ ಚಿಂತೆಮುಕ್ತರಾಗಿ"
    p_guar.font.name = 'Tunga'
    p_guar.font.size = Pt(18)
    p_guar.font.bold = True
    p_guar.font.color.rgb = GREEN
    p_guar.space_after = Pt(10)
    
    p_guar_desc = tf_left8.add_paragraph()
    p_guar_desc.text = "A single complete application remains fully active in the vehicle's engine and lubrication systems for 1,0,000 kilometers, requiring zero top-ups or frequent treatment renewals."
    p_guar_desc.font.name = 'Inter'
    p_guar_desc.font.size = Pt(14)
    p_guar_desc.font.color.rgb = MUTED
    p_guar_desc.space_after = Pt(20)
    
    p_note_header = tf_left8.add_paragraph()
    p_note_header.text = "PRESENTER NOTE:"
    p_note_header.font.name = 'Outfit'
    p_note_header.font.size = Pt(14)
    p_note_header.font.bold = True
    p_note_header.font.color.rgb = CYAN
    p_note_header.space_after = Pt(5)
    
    p_note_text = tf_left8.add_paragraph()
    p_note_text.text = "Keep presentations focused on the long-term cost benefits, mechanical reliability, and peace of mind this 1,0,000 km timeline offers to individual car owners and corporate fleet managers in Karnataka."
    p_note_text.font.name = 'Inter'
    p_note_text.font.size = Pt(13)
    p_note_text.font.color.rgb = MUTED
        
    add_slide_picture(slide8, "longevity_road.png", GREEN)

    # ==========================================
    # SLIDE 9: UNIVERSAL COMPATIBILITY
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_title(slide9, "Compatibility", "Universal Fuel Adaptability", CYAN)
    
    # Left Content
    left_box9 = slide9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
    tf_left9 = left_box9.text_frame
    tf_left9.word_wrap = True
    
    p_intro9 = tf_left9.paragraphs[0]
    p_intro9.text = "Formulated to optimize all standard combustible fuel engine configurations:"
    p_intro9.font.name = 'Inter'
    p_intro9.font.size = Pt(15)
    p_intro9.font.color.rgb = MUTED
    p_intro9.space_after = Pt(25)
    
    bullets9 = [
        ("Petrol Engines (Ethanol Defense):", " Actively guards delicate BS4 & BS6 injectors from corrosive sugarcane-blended E20 fuel deposits."),
        ("Diesel Engines (Heavy Duty):", " Prevents cylinder scaling, exhaust soot buildup, and injectors carbonization under heavy-duty cargo demands."),
        ("Multi-Fuel Performance:", " Works identically on petrol-powered commuter vehicles and diesel-powered cargo fleets, requiring no formula alterations."),
        ("Universal ICE Adaptability:", " Suitable for any Internal Combustion Engine (ICE)—regardless of size, displacement, or application. This includes passenger vehicles, two-wheelers, agricultural tractors, marine outboards, backup generators, and industrial heavy machinery.")
    ]
    
    for title, text in bullets9:
        p_b = tf_left9.add_paragraph()
        p_b.text = "• "
        p_b.font.name = 'Inter'
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = CYAN
        
        run = p_b.add_run()
        run.text = title
        run.font.bold = True
        run.font.color.rgb = WHITE
        
        run2 = p_b.add_run()
        run2.text = text
        run2.font.color.rgb = MUTED
        p_b.space_after = Pt(8)
        
    add_slide_picture(slide9, "universal_vehicles_india.png", CYAN)

    # Save presentation - write to a new visual filename to prevent PowerPoint lock errors
    output_filename = "sundar_carbon_antidote_plus_karnataka_visual.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as {output_filename}")

if __name__ == "__main__":
    create_deck()
